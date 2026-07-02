# -*- coding: utf-8 -*-
"""
Agent Amber v36 - Full autonomous self-learning agent
- Agentic loop: thinks, picks tools, acts, learns
- Self-discovers monday.com board structure
- Web search for market intelligence  
- Google Drive document reading
- Memory via PostgreSQL
- Internal OEP emails sent directly, external go to Paul
- Daily web search, weekly Monday briefing
- Slide decks output as .pptx -> opens directly as Google Slides
  (requires: pip install python-pptx)
"""

import imaplib
import email
import urllib.request
import urllib.parse
import json
import time
import os
import re
import logging
import traceback
import psycopg2
import psycopg2.extras
import hashlib
import hmac
from http.server import HTTPServer, BaseHTTPRequestHandler
import threading
from email.header import decode_header
from datetime import datetime, timezone, timedelta
from anthropic import Anthropic

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
log = logging.getLogger(__name__)

# -- Config --------------------------------------------------------------------
AGENT_EMAIL    = os.environ["AGENT_EMAIL"]
AGENT_PASSWORD = os.environ["AGENT_PASSWORD"]
APPROVER_EMAIL = os.environ["APPROVER_EMAIL"]
SENDGRID_KEY   = os.environ["SENDGRID_API_KEY"]
MONDAY_TOKEN   = os.environ["MONDAY_API_TOKEN"]
DATABASE_URL   = os.environ["DATABASE_URL"]
IMAP_SERVER    = os.environ.get("IMAP_SERVER", "imap.gmail.com")
POLL_INTERVAL  = int(os.environ.get("POLL_INTERVAL", "120"))
AGENT_NAME     = "Agent Amber"
PAUL_EMAIL     = os.environ.get("PAUL_EMAIL", "paul@oceanenergypathway.org")
# -- Access control -------------------------------------------------------------
# Shared drives Amber is allowed to search (exact names, case-insensitive)
ALLOWED_SHARED_DRIVES = [
    "communication and events",
    "countries",
    "fundraising and grants",
    "programs and projects",
    "strategy and policy",
]

# Monday.com board folders/workspaces to block (case-insensitive, partial match)
BLOCKED_MONDAY_FOLDERS = [
    "finance",
    "hr",
    "human resources",
    "salary",
    "salaries",
    "payroll",
    "contracts",
    "funders",
    "fundraising",
]

# Board name keywords that should always be blocked regardless of folder
BLOCKED_BOARD_KEYWORDS = [
    "salary", "salaries", "payroll", "budget", "invoice", "contract",
    "hr ", "human resource", "recruitment", "appraisal", "performance review",
    "funder", "fundrais", "grant application", "donor"
]

def is_board_allowed(board_name: str, workspace_name: str = "") -> bool:
    """Return False if a Monday board should be blocked from Amber."""
    name_lower = board_name.lower()
    workspace_lower = workspace_name.lower()
    # Block if workspace/folder is sensitive
    for blocked in BLOCKED_MONDAY_FOLDERS:
        if blocked in workspace_lower:
            return False
    # Block if board name contains sensitive keywords
    for kw in BLOCKED_BOARD_KEYWORDS:
        if kw in name_lower:
            return False
    return True

def is_drive_allowed(drive_name: str) -> bool:
    """Return True only if a shared drive is in the allowlist."""
    return drive_name.lower().strip() in ALLOWED_SHARED_DRIVES


OEP_DOMAIN     = "oceanenergypathway.org"

# Slack (optional - gracefully disabled if not configured)
SLACK_BOT_TOKEN    = os.environ.get("SLACK_BOT_TOKEN", "")
SLACK_SIGNING_SECRET = os.environ.get("SLACK_SIGNING_SECRET", "")

# Google Drive OAuth (optional - gracefully disabled if not configured)
GOOGLE_CLIENT_ID     = os.environ.get("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = os.environ.get("GOOGLE_CLIENT_SECRET", "")
GOOGLE_REFRESH_TOKEN = os.environ.get("GOOGLE_REFRESH_TOKEN", "")

client = Anthropic()
pending_approvals = {}

# -- Database ------------------------------------------------------------------

def get_db():
    return psycopg2.connect(DATABASE_URL, sslmode="require")

def init_db():
    with get_db() as conn:
        with conn.cursor() as cur:
            cur.execute("""
                CREATE TABLE IF NOT EXISTS amber_memory (
                    id SERIAL PRIMARY KEY,
                    created_at TIMESTAMP DEFAULT NOW(),
                    updated_at TIMESTAMP DEFAULT NOW(),
                    category TEXT NOT NULL,
                    key TEXT NOT NULL,
                    value TEXT NOT NULL,
                    UNIQUE(category, key)
                );
                CREATE TABLE IF NOT EXISTS amber_issues (
                    id SERIAL PRIMARY KEY,
                    first_seen TIMESTAMP DEFAULT NOW(),
                    last_seen TIMESTAMP DEFAULT NOW(),
                    times_seen INT DEFAULT 1,
                    board TEXT,
                    issue_type TEXT,
                    description TEXT,
                    actioned BOOLEAN DEFAULT FALSE
                );
                CREATE TABLE IF NOT EXISTS amber_interactions (
                    id SERIAL PRIMARY KEY,
                    occurred_at TIMESTAMP DEFAULT NOW(),
                    staff_email TEXT,
                    staff_name TEXT,
                    question TEXT,
                    response TEXT
                );
                CREATE TABLE IF NOT EXISTS amber_intel (
                    id SERIAL PRIMARY KEY,
                    captured_at TIMESTAMP DEFAULT NOW(),
                    country TEXT,
                    headline TEXT,
                    summary TEXT,
                    source TEXT,
                    relevance TEXT
                );
            """)
            conn.commit()
    log.info("Database initialised")

def remember(category, key, value):
    """Store something Amber has learned."""
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    INSERT INTO amber_memory (category, key, value, updated_at)
                    VALUES (%s, %s, %s, NOW())
                    ON CONFLICT (category, key) 
                    DO UPDATE SET value = %s, updated_at = NOW()
                """, (category, key, value, value))
                conn.commit()
        return True
    except Exception as e:
        log.error(f"Remember failed: {e}")
        return False

def recall(category, key=None):
    """Read something Amber previously stored."""
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                if key:
                    cur.execute(
                        "SELECT value FROM amber_memory WHERE category = %s AND key = %s",
                        (category, key)
                    )
                    row = cur.fetchone()
                    return row[0] if row else None
                else:
                    cur.execute(
                        "SELECT key, value FROM amber_memory WHERE category = %s ORDER BY updated_at DESC",
                        (category,)
                    )
                    return {row[0]: row[1] for row in cur.fetchall()}
    except Exception as e:
        log.error(f"Recall failed: {e}")
        return None

def track_issue(board, issue_type, description):
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT id, times_seen FROM amber_issues
                    WHERE board = %s AND issue_type = %s AND actioned = FALSE
                    ORDER BY first_seen DESC LIMIT 1
                """, (board, issue_type))
                existing = cur.fetchone()
                if existing:
                    cur.execute(
                        "UPDATE amber_issues SET times_seen = times_seen + 1, last_seen = NOW(), description = %s WHERE id = %s",
                        (description, existing[0])
                    )
                else:
                    cur.execute(
                        "INSERT INTO amber_issues (board, issue_type, description) VALUES (%s, %s, %s)",
                        (board, issue_type, description)
                    )
                conn.commit()
    except Exception as e:
        log.error(f"Track issue failed: {e}")

def get_persistent_issues():
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT board, issue_type, description, times_seen, first_seen
                    FROM amber_issues WHERE times_seen >= 3 AND actioned = FALSE
                    ORDER BY times_seen DESC
                """)
                return cur.fetchall()
    except:
        return []

def save_intel(country, headline, summary, source, relevance):
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO amber_intel (country, headline, summary, source, relevance) VALUES (%s, %s, %s, %s, %s)",
                    (country, headline, summary, source, relevance)
                )
                conn.commit()
    except Exception as e:
        log.error(f"Intel save failed: {e}")

def get_recent_intel(days=7):
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT country, headline, summary, source, captured_at
                    FROM amber_intel
                    WHERE captured_at > NOW() - INTERVAL '%s days'
                    ORDER BY country, captured_at DESC
                """, (days,))
                return cur.fetchall()
    except:
        return []

def save_interaction(staff_email, staff_name, question, response):
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO amber_interactions (staff_email, staff_name, question, response) VALUES (%s, %s, %s, %s)",
                    (staff_email, staff_name, question[:2000], response[:2000])
                )
                conn.commit()
    except Exception as e:
        log.error(f"Interaction save failed: {e}")

def get_staff_history(staff_email, limit=5):
    try:
        with get_db() as conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT question, response, occurred_at FROM amber_interactions
                    WHERE staff_email = %s ORDER BY occurred_at DESC LIMIT %s
                """, (staff_email, limit))
                return cur.fetchall()
    except:
        return []

# -- Monday.com API ------------------------------------------------------------

def monday_api(query):
    data = json.dumps({"query": query}).encode("utf-8")
    auth = f"Bearer {MONDAY_TOKEN}" if not MONDAY_TOKEN.startswith("Bearer") else MONDAY_TOKEN
    req = urllib.request.Request(
        "https://api.monday.com/v2",
        data=data,
        headers={
            "Authorization": auth,
            "Content-Type": "application/json",
            "API-Version": "2024-01"
        },
        method="POST"
    )
    try:
        resp = urllib.request.urlopen(req, timeout=30)
        raw = resp.read()
        result = json.loads(raw)
        # Log if errors returned
        if result.get("errors"):
            log.error(f"Monday API errors: {result['errors']}")
        if result.get("data"):
            boards = result["data"].get("boards", [])
            if boards:
                items = boards[0].get("items_page", {}).get("items", [])
                log.info(f"Monday API returned {len(items)} items from board")
            else:
                log.warning(f"Monday API returned no boards. Full response: {str(result)[:500]}")
        return result
    except urllib.error.HTTPError as e:
        body = e.read().decode("utf-8", errors="replace")
        log.error(f"Monday API HTTP error {e.code}: {body[:500]}")
        raise
    except Exception as e:
        log.error(f"Monday API error: {e}")
        raise

def days_ago(iso_str):
    if not iso_str:
        return None
    try:
        dt = datetime.fromisoformat(iso_str.replace("Z", "+00:00"))
        return (datetime.now(timezone.utc) - dt).days
    except:
        return None

# -- Google Drive --------------------------------------------------------------

_drive_service = None

def get_drive_service():
    global _drive_service
    if _drive_service:
        return _drive_service
    if not GOOGLE_CLIENT_ID or not GOOGLE_REFRESH_TOKEN:
        return None
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request
        import googleapiclient.discovery as discovery

        creds = Credentials(
            token=None,
            refresh_token=GOOGLE_REFRESH_TOKEN,
            client_id=GOOGLE_CLIENT_ID,
            client_secret=GOOGLE_CLIENT_SECRET,
            token_uri="https://oauth2.googleapis.com/token",
            scopes=["https://www.googleapis.com/auth/drive.file", "https://www.googleapis.com/auth/drive.readonly"]
        )
        # Refresh to get a valid access token
        creds.refresh(Request())
        _drive_service = discovery.build("drive", "v3", credentials=creds)
        log.info("Google Drive connected via OAuth")
    except Exception as e:
        log.warning(f"Google Drive not available: {e}")
    return _drive_service


def upload_to_drive(filename, content_bytes, mimetype, folder_name="Amber Infographics"):
    """Upload a file to Google Drive, creating folder if needed. Returns shareable link."""
    token = get_gmail_access_token()  # reuse same OAuth token
    if not token:
        return None, "Google Drive not configured"
    
    try:
        # Find or create folder
        folder_id = None
        search = gmail_api.__func__ if hasattr(gmail_api, '__func__') else None
        
        # Search for existing folder
        import urllib.request, urllib.parse, json
        url = f"https://www.googleapis.com/drive/v3/files?q=name%3D%27{urllib.parse.quote(folder_name)}%27+and+mimeType%3D%27application%2Fvnd.google-apps.folder%27+and+trashed%3Dfalse&fields=files(id,name)"
        req = urllib.request.Request(url, headers={"Authorization": f"Bearer {token}"})
        resp = urllib.request.urlopen(req, timeout=10)
        folders = json.loads(resp.read()).get("files", [])
        
        if folders:
            folder_id = folders[0]["id"]
        else:
            # Create folder
            folder_meta = json.dumps({"name": folder_name, "mimeType": "application/vnd.google-apps.folder"}).encode()
            req = urllib.request.Request(
                "https://www.googleapis.com/drive/v3/files",
                data=folder_meta,
                headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
                method="POST"
            )
            resp = urllib.request.urlopen(req, timeout=10)
            folder_id = json.loads(resp.read())["id"]

        # Multipart upload
        boundary = "amber_boundary_xyz"
        metadata = json.dumps({"name": filename, "parents": [folder_id]}).encode()
        body = (
            f"--{boundary}\r\nContent-Type: application/json\r\n\r\n".encode() +
            metadata + f"\r\n--{boundary}\r\nContent-Type: {mimetype}\r\n\r\n".encode() +
            content_bytes + f"\r\n--{boundary}--".encode()
        )
        req = urllib.request.Request(
            "https://www.googleapis.com/upload/drive/v3/files?uploadType=multipart&fields=id,webViewLink",
            data=body,
            headers={
                "Authorization": f"Bearer {token}",
                "Content-Type": f"multipart/related; boundary={boundary}"
            },
            method="POST"
        )
        resp = urllib.request.urlopen(req, timeout=30)
        result = json.loads(resp.read())
        file_id = result["id"]
        
        # Make it shareable (anyone with link can view)
        perm_data = json.dumps({"role": "reader", "type": "anyone"}).encode()
        req = urllib.request.Request(
            f"https://www.googleapis.com/drive/v3/files/{file_id}/permissions",
            data=perm_data,
            headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
            method="POST"
        )
        urllib.request.urlopen(req, timeout=10)
        
        link = result.get("webViewLink", f"https://drive.google.com/file/d/{file_id}/view")
        return link, None
    except Exception as e:
        log.error(f"Drive upload error: {e}")
        return None, str(e)

# -- Tool definitions for Claude -----------------------------------------------

TOOLS = [
    {
        "name": "explore_board",
        "description": "Explore a monday.com board's structure - discover column names, types, and what data they contain. Use this when you encounter a board you haven't seen before or need to understand its structure. Stores findings for future use.",
        "input_schema": {
            "type": "object",
            "properties": {
                "board_id": {"type": "integer", "description": "The monday.com board ID"},
                "board_name": {"type": "string", "description": "Human-readable name for the board"}
            },
            "required": ["board_id", "board_name"]
        }
    },
    {
        "name": "query_board",
        "description": "Fetch items from a monday.com board. Returns item names, column values, update history, and subitems. Use explore_board first if you haven't seen this board before.",
        "input_schema": {
            "type": "object",
            "properties": {
                "board_id": {"type": "integer", "description": "The monday.com board ID"},
                "limit": {"type": "integer", "description": "Max items to return (default 50)", "default": 50},
                "cursor": {"type": "string", "description": "Pagination cursor for next page"}
            },
            "required": ["board_id"]
        }
    },
    {
        "name": "get_activity",
        "description": "Get the activity log for a monday.com board. Captures ALL types of activity including comments, status changes, file uploads, new items - anything that happened on the board.",
        "input_schema": {
            "type": "object",
            "properties": {
                "board_id": {"type": "integer", "description": "The monday.com board ID"},
                "limit": {"type": "integer", "description": "Number of activity events to return (default 50)", "default": 50}
            },
            "required": ["board_id"]
        }
    },
    {
        "name": "web_search",
        "description": "Search the internet for current information. Use for market intelligence, regulatory updates, news about offshore wind in specific countries, industry developments.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query"},
                "context": {"type": "string", "description": "Why you're searching - helps focus results"}
            },
            "required": ["query"]
        }
    },
    {
        "name": "search_drive",
        "description": "Search OEP's Google Drive for relevant documents. Returns file names, types, and snippets.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "What to search for"},
                "file_type": {"type": "string", "description": "Optional: 'doc', 'pdf', 'sheet', 'presentation'"}
            },
            "required": ["query"]
        }
    },
    {
        "name": "read_document",
        "description": "Read the content of a specific Google Drive document by its file ID.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_id": {"type": "string", "description": "Google Drive file ID"},
                "file_name": {"type": "string", "description": "File name for context"}
            },
            "required": ["file_id"]
        }
    },
    {
        "name": "remember",
        "description": "Store something you've learned for future use. Use this to save board structure discoveries, key facts about OEP's programme, patterns you've noticed.",
        "input_schema": {
            "type": "object",
            "properties": {
                "category": {"type": "string", "description": "Category e.g. 'board_structure', 'country_context', 'staff', 'market_intel'"},
                "key": {"type": "string", "description": "Unique key within category"},
                "value": {"type": "string", "description": "What to remember"}
            },
            "required": ["category", "key", "value"]
        }
    },
    {
        "name": "recall",
        "description": "Read something you previously stored. Use this before querying boards to check if you already know their structure.",
        "input_schema": {
            "type": "object",
            "properties": {
                "category": {"type": "string", "description": "Category to recall from"},
                "key": {"type": "string", "description": "Specific key, or omit to get all entries in category"}
            },
            "required": ["category"]
        }
    },
    {
        "name": "flag_issue",
        "description": "Flag a recurring issue that needs attention. Issues flagged 3+ times in a row trigger escalation in the Monday briefing.",
        "input_schema": {
            "type": "object",
            "properties": {
                "board": {"type": "string", "description": "Which board/country this relates to"},
                "issue_type": {"type": "string", "description": "Type of issue e.g. 'stale_items', 'missing_owners', 'data_quality'"},
                "description": {"type": "string", "description": "Clear description of the issue"}
            },
            "required": ["board", "issue_type", "description"]
        }
    },
    {
        "name": "create_infographic",
        "description": "Create a visual slide deck (.pptx) from data and save it to Google Drive. The file opens directly as Google Slides. Use when asked to create visual summaries, status dashboards, project overviews, slide decks, briefings, or any data visualisation. You provide the title and sections - Amber generates a clean styled presentation and returns a shareable Drive link that opens in Google Slides.",
        "input_schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Title of the presentation"},
                "subtitle": {"type": "string", "description": "Optional subtitle or date"},
                "sections": {
                    "type": "array",
                    "description": "Sections of content - each becomes one or more slides",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "type": {"type": "string", "enum": ["text", "table", "status_grid", "stat_row", "bullets"]},
                            "content": {"type": "string", "description": "Text content, or JSON string for structured data"}
                        }
                    }
                },
                "colour_scheme": {"type": "string", "enum": ["ocean", "corporate", "warm", "dark"], "default": "ocean"}
            },
            "required": ["title", "sections"]
        }
    },
    {
        "name": "read_paul_inbox",
        "description": "Read Paul's email inbox and return a summary of recent emails. ONLY available when responding to Paul (paul@oceanenergypathway.org) - never use this for any other sender. Use when Paul asks about his emails, wants a summary, asks what he's missed, or asks you to remind him of things.",
        "input_schema": {
            "type": "object",
            "properties": {
                "hours": {"type": "integer", "description": "How many hours back to look (default 24)", "default": 24},
                "focus": {"type": "string", "description": "Optional: what to focus on e.g. 'unread only', 'from a specific person'"}
            },
            "required": []
        }
    },
    {
        "name": "read_slack_channel",
        "description": "Read recent messages from a Slack channel Amber has been invited to. Use this to understand what the team is discussing, find context about projects, or answer questions about conversations. Returns the last N messages with sender names and timestamps.",
        "input_schema": {
            "type": "object",
            "properties": {
                "channel_name": {"type": "string", "description": "Channel name e.g. '#general', '#japan', '#programme' - or 'all' to get a summary across all channels"},
                "hours": {"type": "integer", "description": "How many hours back to look (default 48)", "default": 48},
                "limit": {"type": "integer", "description": "Max messages to return (default 30)", "default": 30}
            },
            "required": ["channel_name"]
        }
    },
    {
        "name": "finish",
        "description": "Call this when you are ready to send your response. The 'response' field IS the email that gets sent immediately. Write the COMPLETE email content here - do not summarise what you are about to write, do not say 'let me compile this now', do not use this as a placeholder. The full briefing, analysis, or answer goes directly in the response field. This is the last tool you call.",
        "input_schema": {
            "type": "object",
            "properties": {
                "response": {"type": "string", "description": "The complete email content to send. Must be the full response, not a placeholder or summary of what you will write."}
            },
            "required": ["response"]
        }
    }
]

# -- Tool execution ------------------------------------------------------------

def execute_tool(tool_name, tool_input):
    """Execute a tool call and return the result."""
    
    if tool_name == "explore_board":
        board_id = tool_input["board_id"]
        board_name = tool_input["board_name"]
        try:
            # First check if we already know this board
            cached = recall("board_structure", str(board_id))
            if cached:
                return f"Already know this board (from memory): {cached}"
            
            result = monday_api(f"""
            {{
              boards(ids: [{board_id}]) {{
                name
                description
                columns {{
                  id
                  title
                  type
                  description
                }}
                groups {{
                  id
                  title
                }}
                items_page(limit: 3) {{
                  items {{
                    name
                    column_values {{ id text }}
                  }}
                }}
              }}
            }}
            """)
            
            board = result.get("data", {}).get("boards", [{}])[0]
            columns = board.get("columns", [])
            sample_items = board.get("items_page", {}).get("items", [])
            
            # Build a clear picture of the board structure
            col_summary = []
            for col in columns:
                sample_vals = []
                for item in sample_items:
                    for cv in item.get("column_values", []):
                        if cv["id"] == col["id"] and cv.get("text"):
                            sample_vals.append(cv["text"])
                col_summary.append({
                    "id": col["id"],
                    "title": col["title"],
                    "type": col["type"],
                    "sample_values": sample_vals[:2]
                })
            
            structure = json.dumps(col_summary, indent=2)
            
            # Store for future use
            remember("board_structure", str(board_id), structure)
            remember("board_structure", f"{board_id}_name", board_name)
            
            return f"Board '{board_name}' structure discovered and saved:\n{structure}"
            
        except Exception as e:
            return f"Error exploring board {board_id}: {e}"

    elif tool_name == "query_board":
        board_id = tool_input["board_id"]
        limit = tool_input.get("limit", 50)
        cursor = tool_input.get("cursor", "")
        cursor_part = f', cursor: "{cursor}"' if cursor else ""
        
        try:
            result = monday_api(f"""
            {{
              boards(ids: [{board_id}]) {{
                name
                items_page(limit: {limit}{cursor_part}) {{
                  cursor
                  items {{
                    id
                    name
                    state
                    created_at
                    updated_at
                    creator {{ name }}
                    column_values {{ id text }}
                    updates(limit: 2) {{
                      body
                      created_at
                      creator {{ name }}
                    }}
                    subitems {{
                      id
                      name
                      updated_at
                      column_values {{ id text }}
                    }}
                  }}
                }}
              }}
            }}
            """)

            
            board = result.get("data", {}).get("boards", [{}])[0]
            items = board.get("items_page", {}).get("items", [])
            next_cursor = board.get("items_page", {}).get("cursor")
            
            # Enrich items with computed fields
            enriched = []
            for item in items:
                if item.get("state") == "deleted":
                    continue
                
                # Build column map by both ID and title
                cols = {}
                for cv in item.get("column_values", []):
                    if cv.get("text"):
                        cols[cv["id"]] = cv["text"]
                        cols[cv.get("title", "").lower()] = cv["text"]
                
                enriched.append({
                    "id": item["id"],
                    "name": item["name"],
                    "created_days_ago": days_ago(item.get("created_at")),
                    "updated_days_ago": days_ago(item.get("updated_at")),
                    "creator": item.get("creator", {}).get("name", ""),
                    "columns": cols,
                    "latest_updates": [
                        {
                            "body": u.get("body", "")[:200],
                            "by": u.get("creator", {}).get("name", ""),
                            "days_ago": days_ago(u.get("created_at"))
                        }
                        for u in item.get("updates", [])
                    ],
                    "subitems": [
                        {
                            "name": s["name"],
                            "updated_days_ago": days_ago(s.get("updated_at")),
                            "status": next(
                                (cv["text"] for cv in s.get("column_values", []) if cv.get("text")),
                                "not set"
                            )
                        }
                        for s in item.get("subitems", [])
                    ]
                })
            
            return json.dumps({
                "board_name": board.get("name"),
                "item_count": len(enriched),
                "has_more": bool(next_cursor),
                "next_cursor": next_cursor,
                "items": enriched
            }, default=str)
            
        except Exception as e:
            return f"Error querying board {board_id}: {e}\n{traceback.format_exc()}"

    elif tool_name == "get_activity":
        board_id = tool_input["board_id"]
        limit = tool_input.get("limit", 50)
        try:
            result = monday_api(f"""
            {{
              boards(ids: [{board_id}]) {{
                activity_logs(limit: {limit}) {{
                  id
                  event
                  created_at
                  data
                }}
              }}
            }}
            """)
            
            logs = result.get("data", {}).get("boards", [{}])[0].get("activity_logs", [])
            now = datetime.now(timezone.utc)
            
            processed = []
            for entry in logs:
                ts = entry.get("created_at")
                d = None
                if ts:
                    try:
                        ts_sec = int(ts) / 10_000_000
                        dt = datetime.fromtimestamp(ts_sec, tz=timezone.utc)
                        d = (now - dt).days
                    except:
                        pass
                processed.append({
                    "event": entry.get("event", ""),
                    "by": entry.get("user", {}).get("name", "system") if entry.get("user") else "system",
                    "days_ago": d,
                    "data": entry.get("data", "")[:100]
                })
            
            return json.dumps({"activity_count": len(processed), "activity": processed}, default=str)
            
        except Exception as e:
            return f"Error getting activity for board {board_id}: {e}"

    elif tool_name == "web_search":
        query = tool_input["query"]
        try:
            response = client.messages.create(
                model="claude-sonnet-4-6",
                max_tokens=1000,
                tools=[{"type": "web_search_20250305", "name": "web_search"}],
                messages=[{"role": "user", "content": f"Search for: {query}. Return a clear summary of the most relevant and recent findings."}]
            )
            text = " ".join(block.text for block in response.content if hasattr(block, "text"))
            
            # Save relevant intel
            save_intel(query, query[:100], text[:500], "web_search", tool_input.get("context", ""))
            
            return text[:2000]
        except Exception as e:
            return f"Web search error: {e}"

    elif tool_name == "search_drive":
        drive = get_drive_service()
        if not drive:
            return "Google Drive not connected. Add GOOGLE_CREDENTIALS to Railway variables to enable."
        try:
            query = tool_input["query"]
            file_type = tool_input.get("file_type", "")

            # Get allowed shared drive IDs
            all_drives = drive.drives().list(fields="drives(id,name)").execute()
            allowed_drive_ids = [
                d["id"] for d in all_drives.get("drives", [])
                if is_drive_allowed(d["name"])
            ]

            all_files = []
            for drive_id in allowed_drive_ids:
                q = f"fullText contains '{query}' and trashed = false"
                if file_type == "doc":
                    q += " and mimeType = 'application/vnd.google-apps.document'"
                elif file_type == "sheet":
                    q += " and mimeType = 'application/vnd.google-apps.spreadsheet'"
                elif file_type == "pdf":
                    q += " and mimeType = 'application/pdf'"
                try:
                    results = drive.files().list(
                        q=q,
                        pageSize=8,
                        corpora="drive",
                        driveId=drive_id,
                        includeItemsFromAllDrives=True,
                        supportsAllDrives=True,
                        fields="files(id, name, mimeType, modifiedTime, driveId)"
                    ).execute()
                    all_files.extend(results.get("files", []))
                except Exception:
                    continue

            # Deduplicate by id
            seen = set()
            unique_files = []
            for f in all_files:
                if f["id"] not in seen:
                    seen.add(f["id"])
                    unique_files.append(f)

            return json.dumps([{
                "id": f["id"],
                "name": f["name"],
                "type": f["mimeType"],
                "modified": f.get("modifiedTime", "")
            } for f in unique_files[:10]])

        except Exception as e:
            return f"Drive search error: {e}"

    elif tool_name == "read_document":
        drive = get_drive_service()
        if not drive:
            return "Google Drive not connected."
        try:
            file_id = tool_input["file_id"]
            
            # Export as plain text
            content = drive.files().export(
                fileId=file_id,
                mimeType="text/plain"
            ).execute()
            
            if isinstance(content, bytes):
                content = content.decode("utf-8", errors="replace")
            
            return content[:5000]  # First 5000 chars
            
        except Exception as e:
            return f"Document read error: {e}"

    elif tool_name == "remember":
        success = remember(tool_input["category"], tool_input["key"], tool_input["value"])
        return "Saved to memory." if success else "Memory save failed."

    elif tool_name == "recall":
        result = recall(tool_input["category"], tool_input.get("key"))
        if result is None:
            return "Nothing found in memory for that category/key."
        return json.dumps(result) if isinstance(result, dict) else str(result)

    elif tool_name == "flag_issue":
        track_issue(tool_input["board"], tool_input["issue_type"], tool_input["description"])
        return f"Issue flagged: {tool_input['board']} - {tool_input['description']}"

    elif tool_name == "create_infographic":
        title = tool_input.get("title", "OEP Report")
        subtitle = tool_input.get("subtitle", "")
        sections = tool_input.get("sections", [])
        scheme = tool_input.get("colour_scheme", "ocean")

        # Colour palettes - (primary_hex, accent_hex, text_on_primary_hex)
        palettes = {
            "ocean":     {"primary": (0x00, 0x77, 0xB6), "accent": (0x90, 0xE0, 0xEF), "dark": (0x03, 0x04, 0x5E)},
            "corporate": {"primary": (0x2C, 0x3E, 0x50), "accent": (0x34, 0x98, 0xDB), "dark": (0x1a, 0x1a, 0x2e)},
            "warm":      {"primary": (0xE7, 0x6F, 0x51), "accent": (0xF4, 0xA2, 0x61), "dark": (0x3D, 0x40, 0x5B)},
            "dark":      {"primary": (0x6C, 0x63, 0xFF), "accent": (0x3F, 0x3D, 0x56), "dark": (0x1a, 0x1a, 0x2e)},
        }
        pal = palettes.get(scheme, palettes["ocean"])

        try:
            import io
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt

            PRIMARY = RGBColor(*pal["primary"])
            ACCENT  = RGBColor(*pal["accent"])
            DARK    = RGBColor(*pal["dark"])
            WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
            LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)
            MID     = RGBColor(0x44, 0x44, 0x44)

            SW = Inches(13.33)  # widescreen 16:9
            SH = Inches(7.5)

            prs = Presentation()
            prs.slide_width  = SW
            prs.slide_height = SH

            def blank_slide(prs):
                layout = prs.slide_layouts[6]  # completely blank
                return prs.slides.add_slide(layout)

            def fill_shape(shape, rgb):
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb

            def add_rect(slide, left, top, width, height, rgb):
                shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
                fill_shape(shape, rgb)
                shape.line.fill.background()
                return shape

            def add_textbox(slide, text, left, top, width, height,
                            font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
                txb = slide.shapes.add_textbox(left, top, width, height)
                tf  = txb.text_frame
                tf.word_wrap = wrap
                p   = tf.paragraphs[0]
                p.alignment = align
                run = p.add_run()
                run.text = str(text)
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = color or DARK
                return txb

            # -- TITLE SLIDE --------------------------------------------------
            sl = blank_slide(prs)
            add_rect(sl, 0, 0, SW, SH, PRIMARY)                          # full bg
            add_rect(sl, 0, SH - Inches(0.08), SW, Inches(0.08), ACCENT) # bottom strip

            # Title
            add_textbox(sl, title,
                        Inches(1), Inches(2.2), Inches(11.33), Inches(1.8),
                        font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            # Subtitle
            if subtitle:
                add_textbox(sl, subtitle,
                            Inches(1), Inches(4.1), Inches(11.33), Inches(0.7),
                            font_size=20, color=ACCENT, align=PP_ALIGN.CENTER)
            # Footer
            generated = datetime.now(timezone.utc).strftime("%d %b %Y")
            add_textbox(sl, f"Generated by Agent Amber * {generated}",
                        Inches(1), Inches(6.8), Inches(11.33), Inches(0.4),
                        font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

            # -- CONTENT SLIDES -----------------------------------------------
            for sec in sections:
                heading     = sec.get("heading", "")
                stype       = sec.get("type", "text")
                content_raw = sec.get("content", "")
                if not isinstance(content_raw, str):
                    content_raw = json.dumps(content_raw)

                sl = blank_slide(prs)

                # Header bar
                add_rect(sl, 0, 0, SW, Inches(1.1), PRIMARY)
                add_rect(sl, 0, Inches(1.1), SW, Inches(0.05), ACCENT)

                # Slide title in header
                add_textbox(sl, heading,
                            Inches(0.5), Inches(0.15), Inches(12.33), Inches(0.8),
                            font_size=26, bold=True, color=WHITE)

                # Footer strip
                add_rect(sl, 0, SH - Inches(0.35), SW, Inches(0.35), LIGHT)
                add_textbox(sl, title,
                            Inches(0.4), SH - Inches(0.32), Inches(8), Inches(0.28),
                            font_size=9, color=MID)

                CONTENT_TOP  = Inches(1.3)
                CONTENT_H    = Inches(5.9)
                CONTENT_L    = Inches(0.5)
                CONTENT_W    = Inches(12.33)

                if stype == "bullets":
                    # Parse bullet list
                    if isinstance(content_raw, list):
                        lines = content_raw
                    else:
                        try:
                            parsed = json.loads(content_raw)
                            lines = parsed if isinstance(parsed, list) else content_raw.split("\n")
                        except:
                            lines = content_raw.split("\n")
                    lines = [str(l).lstrip("*-* ").strip() for l in lines if str(l).strip()]

                    txb = slide.shapes.add_textbox if False else sl.shapes.add_textbox(
                        CONTENT_L, CONTENT_TOP, CONTENT_W, CONTENT_H)
                    tf = txb.text_frame
                    tf.word_wrap = True
                    for i, line in enumerate(lines[:12]):
                        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                        p.text = f"  *  {line}"
                        p.space_before = Pt(6)
                        run = p.runs[0] if p.runs else p.add_run()
                        run.font.size = Pt(18)
                        run.font.color.rgb = DARK

                elif stype == "stat_row":
                    try:
                        stats = json.loads(content_raw)
                    except:
                        stats = []
                    n = max(len(stats), 1)
                    box_w = Inches(12.33 / n) - Inches(0.15)
                    for i, s in enumerate(stats[:6]):
                        lft = Inches(0.5) + i * (box_w + Inches(0.15))
                        # Box background
                        add_rect(sl, lft, CONTENT_TOP + Inches(0.3), box_w, Inches(2.2), LIGHT)
                        # Value
                        add_textbox(sl, str(s.get("value", "")),
                                    lft, CONTENT_TOP + Inches(0.4), box_w, Inches(1.4),
                                    font_size=40, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
                        # Label
                        add_textbox(sl, str(s.get("label", "")),
                                    lft, CONTENT_TOP + Inches(1.8), box_w, Inches(0.6),
                                    font_size=13, color=MID, align=PP_ALIGN.CENTER)

                elif stype == "status_grid":
                    try:
                        items = json.loads(content_raw)
                    except:
                        items = []
                    STATUS_COLORS = {
                        "on-track":  (RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x1B, 0x5E, 0x20)),
                        "on track":  (RGBColor(0xE8, 0xF5, 0xE9), RGBColor(0x1B, 0x5E, 0x20)),
                        "at-risk":   (RGBColor(0xFF, 0xF9, 0xC4), RGBColor(0x7A, 0x5C, 0x00)),
                        "at risk":   (RGBColor(0xFF, 0xF9, 0xC4), RGBColor(0x7A, 0x5C, 0x00)),
                        "blocked":   (RGBColor(0xFF, 0xEB, 0xEE), RGBColor(0x8B, 0x1A, 0x1A)),
                        "done":      (RGBColor(0xE8, 0xEA, 0xFF), RGBColor(0x1A, 0x3A, 0x8B)),
                        "neutral":   (RGBColor(0xF5, 0xF5, 0xF5), MID),
                    }
                    cols = 4
                    box_w = Inches(2.9)
                    box_h = Inches(0.85)
                    gap   = Inches(0.15)
                    for idx, item in enumerate(items[:16]):
                        col = idx % cols
                        row = idx // cols
                        lft = Inches(0.5) + col * (box_w + gap)
                        top = CONTENT_TOP + Inches(0.1) + row * (box_h + gap)
                        status_key = item.get("status", "neutral").lower()
                        bg, fg = STATUS_COLORS.get(status_key, STATUS_COLORS["neutral"])
                        add_rect(sl, lft, top, box_w, box_h, bg)
                        add_textbox(sl, item.get("name", ""),
                                    lft + Inches(0.1), top + Inches(0.05), box_w - Inches(0.2), Inches(0.4),
                                    font_size=13, bold=True, color=fg)
                        add_textbox(sl, item.get("status", ""),
                                    lft + Inches(0.1), top + Inches(0.45), box_w - Inches(0.2), Inches(0.3),
                                    font_size=11, color=fg)

                elif stype == "table":
                    try:
                        rows = json.loads(content_raw)
                    except:
                        rows = []
                    if rows:
                        headers = list(rows[0].keys())
                        ncols = len(headers)
                        nrows = min(len(rows), 10)
                        col_w = [Inches(12.33 / ncols)] * ncols
                        row_h = [Inches(0.45)] * (nrows + 1)
                        tbl = sl.shapes.add_table(nrows + 1, ncols,
                                                   CONTENT_L, CONTENT_TOP,
                                                   Inches(12.33), Inches(0.45 * (nrows + 1))).table
                        tbl.columns
                        for ci, h in enumerate(headers):
                            cell = tbl.cell(0, ci)
                            cell.text = h.upper()
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = PRIMARY
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.color.rgb = WHITE
                                    run.font.bold = True
                                    run.font.size = Pt(12)
                        for ri, row in enumerate(rows[:10]):
                            for ci, h in enumerate(headers):
                                cell = tbl.cell(ri + 1, ci)
                                cell.text = str(row.get(h, ""))
                                if ri % 2 == 0:
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = LIGHT
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = DARK

                else:  # plain text
                    add_textbox(sl, content_raw,
                                CONTENT_L, CONTENT_TOP, CONTENT_W, CONTENT_H,
                                font_size=18, color=DARK)

            # -- Save to bytes ------------------------------------------------
            buf = io.BytesIO()
            prs.save(buf)
            pptx_bytes = buf.getvalue()

        except ImportError:
            return ("python-pptx is not installed. Run: pip install python-pptx\n"
                    "Then redeploy the agent.")
        except Exception as e:
            log.error(f"PPTX generation error: {e}\n{traceback.format_exc()}")
            return f"Slide deck generation failed: {e}"

        # -- Upload as .pptx - Google Drive auto-converts to Google Slides --
        safe_name = re.sub(r"[^\w\s-]", "", title).strip().replace(" ", "_")[:50]
        filename  = safe_name + ".pptx"
        mime      = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        link, err = upload_to_drive(filename, pptx_bytes, mime, folder_name="Amber Presentations")
        if err:
            return f"Slide deck created but could not upload to Drive: {err}"

        # -- Also email as attachment -----------------------------------------
        try:
            import base64, email.mime.multipart, email.mime.text, email.mime.base, email.encoders
            token = get_gmail_access_token()
            if token:
                msg = email.mime.multipart.MIMEMultipart()
                msg["to"] = PAUL_EMAIL
                msg["subject"] = f"[Amber] {title}"
                body_part = email.mime.text.MIMEText(
                    f"Hi Paul,\n\nYour slide deck is ready.\n\n"
                    f"Open in Google Slides: {link}\n\n"
                    f"(The .pptx attachment also opens in Google Slides or PowerPoint.)\n\n- Amber",
                    "plain"
                )
                msg.attach(body_part)
                attach = email.mime.base.MIMEBase(
                    "application",
                    "vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                attach.set_payload(pptx_bytes)
                email.encoders.encode_base64(attach)
                attach.add_header("Content-Disposition", "attachment", filename=filename)
                msg.attach(attach)
                raw = base64.urlsafe_b64encode(msg.as_bytes()).decode()
                send_req = urllib.request.Request(
                    "https://gmail.googleapis.com/gmail/v1/users/me/messages/send",
                    data=json.dumps({"raw": raw}).encode(),
                    headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
                    method="POST"
                )
                urllib.request.urlopen(send_req, timeout=15)
        except Exception as e:
            log.error(f"Slide deck email failed: {e}")

        return (
            f"Slide deck created ({len(sections)} slides)!\n\n"
            f"Open in Google Slides: {link}\n\n"
            f"I've also emailed it to you as a .pptx - click the Drive link and it will open directly in Google Slides."
        )


    elif tool_name == "read_paul_inbox":
        # Security gate - only Paul can use this tool
        sender_email = tool_input.get("_sender_email", "")
        if sender_email.lower() != PAUL_EMAIL.lower():
            return "Access denied - inbox reading is only available to Paul."
        hours = tool_input.get("hours", 24)
        return read_paul_inbox(hours=hours, max_emails=40)


    elif tool_name == "read_slack_channel":
        if not SLACK_BOT_TOKEN:
            return "Slack is not configured."
        channel_name = tool_input.get("channel_name", "all")
        hours = tool_input.get("hours", 48)
        limit = tool_input.get("limit", 30)
        oldest = str((datetime.now(timezone.utc) - timedelta(hours=hours)).timestamp())

        # Get bot user ID for name resolution
        auth = slack_get("auth.test", {})
        bot_user_id = auth.get("user_id", "") if auth else ""

        # Build user ID -> name cache
        users_result = slack_get("users.list", {"limit": 200})
        user_map = {}
        if users_result and users_result.get("ok"):
            for u in users_result.get("members", []):
                user_map[u["id"]] = u.get("real_name") or u.get("name", u["id"])

        # Get channels
        channels_result = slack_get("conversations.list", {
            "types": "public_channel",
            "exclude_archived": "true",
            "limit": 100
        })
        if not channels_result or not channels_result.get("ok"):
            return "Could not retrieve Slack channels."

        channels = [c for c in channels_result.get("channels", []) if c.get("is_member")]

        # Filter to requested channel or all
        if channel_name != "all":
            clean = channel_name.lstrip("#").lower()
            channels = [c for c in channels if c["name"].lower() == clean]
            if not channels:
                return f"Channel #{clean} not found or Amber hasn't been invited to it."

        all_messages = []
        for ch in channels:
            result = slack_get("conversations.history", {
                "channel": ch["id"],
                "oldest": oldest,
                "limit": limit
            })
            if not result or not result.get("ok"):
                continue
            for msg in result.get("messages", []):
                if msg.get("bot_id") or msg.get("subtype"):
                    continue
                user_id = msg.get("user", "")
                if user_id == bot_user_id:
                    continue
                name = user_map.get(user_id, user_id)
                ts = float(msg.get("ts", 0))
                dt = datetime.fromtimestamp(ts, tz=timezone.utc).strftime("%a %d %b %H:%M")
                text = msg.get("text", "").strip()
                if text:
                    all_messages.append({
                        "channel": ch["name"],
                        "time": dt,
                        "name": name,
                        "text": text[:300]
                    })

        if not all_messages:
            return f"No messages found in the last {hours} hours."

        # Sort by time
        all_messages.sort(key=lambda x: x["time"])

        lines = [f"[#{m['channel']}] {m['time']} - {m['name']}: {m['text']}" for m in all_messages]
        return f"Found {len(lines)} messages:\n\n" + "\n".join(lines)


    elif tool_name == "finish":
        return "__FINISH__"

    return f"Unknown tool: {tool_name}"

# -- Agentic loop --------------------------------------------------------------

SYSTEM_PROMPT = """
You are Agent Amber, OEP's Programme Intelligence Agent for Ocean Energy Pathway.
OEP is an independent non-profit accelerating offshore wind in 9 emerging markets.
You are professional, direct, warm, and specific. Sign emails as "Agent Amber".

You have tools to explore monday.com boards, search the web, read Google Drive documents,
and store/retrieve your own memories. Use them intelligently.

CORE PRINCIPLES:
1. Always recall memory before querying - you may already know the answer
2. Always explore_board before query_board if you haven't seen it before
3. Be specific - name items, people, dates. Never generalise.
4. Skip completed/closed projects in operational analysis
5. Cross-reference master board (ID: 1747605081) with country boards
6. Flag issues that recur - they need escalation

OEP BOARD STRUCTURE:
- Master board ID: 1747605081 - "OEP Projects (DO NOT EDIT HERE)" - source of truth for status, timeline, owner, budget, funder
- Country boards mirror from master - query master for governance data, country boards for activity/comments
- Countries: Brazil (1767248587), India (1767246703), Japan (1767245536), Philippines (1767246398), South Korea (1767190694), Vietnam (2053944026), Mexico (2007096589), Australia (1955282782), Colombia (1879431534)
- Global projects (1767247726) - skip items where Location = QA

PROJECT CLOSURE RULES:
- Closed = all subitems marked Done or Skip
- Effectively closed = only "Final Evaluation with Chidinma" remaining open
- If status on master board = "Completed" -> skip from operational analysis
- Monitoring items = description contains "Monitoring item - MEL final stages to be completed"

MARKET PRIORITIES (staleness thresholds):
- Japan, South Korea: 30 days (flagship - advisory roles active)
- India, Brazil, Philippines: 45 days (high priority)
- Vietnam, Colombia, Australia: 60 days (medium - note Australia auction Aug 2026 is IMMINENT)
- Mexico: 180 days (early stage, 3 projects is normal)

OKR REPORTING:
- Q1=Jan-Mar, Q2=Apr-Jun, Q3=Jul-Sep, Q4=Oct-Dec
- OKR 1.3: 90% on schedule - needs timelines set
- Effectively closed projects count as delivered for that quarter
- Count dissemination events: public events, publications, govt presentations, stakeholder workshops

DATA QUALITY - flag per item:
- No status set -> blocks progress tracking
- No timeline -> blocks OKR 1.3
- No owner -> no accountability
- Never updated since creation -> likely placeholder

TONE: Professional, collegial, specific. Avoid dramatic language.
Always include a "Board improvements" section in full briefings.

When you have enough information, call the finish tool with your COMPLETE response.
CRITICAL: The text you put in the finish tool IS the email that gets sent. 
Do NOT say "I now have everything I need" or "Let me write the briefing" - write the actual briefing IN the finish tool.
Do NOT use finish as a placeholder. The finish tool content = the email content. Write it all there.
"""

def run_agentic_loop(question, sender_name, sender_email, max_iterations=25):
    """
    Run the agentic loop: Amber thinks, picks tools, acts, learns, repeats.
    Returns the final response string.

    Conversation rules enforced here:
      - Every tool_use block MUST be followed by a tool_result in the very next
        user message. No plain text user messages may be injected while
        tool_results are pending.
      - Nudges (force-finish, spiral warnings) are piggybacked onto the content
        of the last real tool_result so the conversation structure stays valid.
    """
    log.info(f"Starting agentic loop for: {sender_email}")

    history      = get_staff_history(sender_email)
    persistent   = get_persistent_issues()
    recent_intel = get_recent_intel(days=7)

    context_parts = [f"Email from: {sender_name} <{sender_email}>\n\n{question}"]

    if history:
        context_parts.append("\nPREVIOUS INTERACTIONS WITH THIS PERSON:")
        for q, r, ts in history[:3]:
            context_parts.append(f"[{ts.strftime('%d %b %Y')}] Asked: {q[:150]}")

    if persistent:
        context_parts.append("\nPERSISTENT ISSUES (flagged 3+ consecutive checks):")
        for board, issue_type, desc, times, first_seen in persistent:
            context_parts.append(f"- {board}: {desc} (seen {times} times since {first_seen.strftime('%d %b')})")

    if recent_intel:
        context_parts.append("\nRECENT MARKET INTELLIGENCE (last 7 days):")
        for country, headline, summary, source, ts in recent_intel[:6]:
            context_parts.append(f"- {country}: {headline}")

    messages = [{"role": "user", "content": "\n".join(context_parts)}]
    recent_tool_calls = []

    for iteration in range(max_iterations):
        log.info(f"Agentic loop iteration {iteration + 1}")

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=4096,
            system=SYSTEM_PROMPT,
            tools=TOOLS,
            messages=messages
        )

        messages.append({"role": "assistant", "content": response.content})

        # end_turn means no tool calls — extract text and return
        if response.stop_reason == "end_turn":
            text = " ".join(
                block.text for block in response.content
                if hasattr(block, "text") and block.text.strip()
            )
            if text.strip():
                return text
            break

        # Collect ALL tool_use blocks before executing any
        tool_calls = [b for b in response.content if b.type == "tool_use"]
        if not tool_calls:
            break

        tool_results   = []
        final_response = None
        iteration_tools = []

        for idx, block in enumerate(tool_calls):
            tool_name  = block.name
            tool_input = dict(block.input)
            log.info(f"Tool call: {tool_name}({list(tool_input.keys())})")
            iteration_tools.append(tool_name)
            tool_input["_sender_email"] = sender_email
            result = execute_tool(tool_name, tool_input)

            if result == "__FINISH__":
                final_response = tool_input.get("response", "")
                tool_results.append({
                    "type": "tool_result",
                    "tool_use_id": block.id,
                    "content": "done"
                })
                # Stub all remaining tool calls in this batch
                for rem in tool_calls[idx + 1:]:
                    log.info(f"Stub result for unprocessed tool: {rem.name}")
                    tool_results.append({
                        "type": "tool_result",
                        "tool_use_id": rem.id,
                        "content": "skipped - finish was called"
                    })
                break

            tool_results.append({
                "type": "tool_result",
                "tool_use_id": block.id,
                "content": str(result)[:8000]
            })

        # Post all tool results before returning or looping
        if tool_results:
            # Piggyback any nudge onto the last result's content string
            # so we never inject a bare user text message mid-conversation.
            recent_tool_calls.extend(iteration_tools)
            recent_tool_calls = recent_tool_calls[-6:]

            nudge = None
            if iteration == max_iterations - 4:
                nudge = (
                    "[SYSTEM NOTE: You are approaching the iteration limit. "
                    "Please call the finish tool on your next response with "
                    "everything you have found. Write a complete reply.]"
                )
            elif (len(recent_tool_calls) == 6 and
                    all(t in ("search_drive", "read_document", "web_search")
                        for t in recent_tool_calls)):
                log.warning("Search spiral detected - piggybacking nudge onto tool results")
                nudge = (
                    "[SYSTEM NOTE: You appear stuck in a search loop. "
                    "Stop searching and call the finish tool now with what you have found. "
                    "If you could not find everything, say so clearly.]"
                )
                recent_tool_calls = []

            if nudge and tool_results:
                tool_results[-1]["content"] = str(tool_results[-1]["content"]) + "\n\n" + nudge

            messages.append({"role": "user", "content": tool_results})

        if final_response is not None:
            if final_response.strip():
                return final_response
            log.warning("finish tool called with empty response - continuing loop")

        if not tool_results:
            break

    # Fallback: return any text from the last assistant message
    for msg in reversed(messages):
        if msg.get("role") == "assistant":
            content = msg.get("content", [])
            if isinstance(content, list):
                text = " ".join(
                    block.text for block in content
                    if hasattr(block, "text") and block.text.strip()
                )
                if text.strip():
                    log.warning("Returning fallback text from last assistant message")
                    return text
            break

    return (
        "I ran into difficulties completing this analysis - "
        "I may have been unable to access the files or data needed. "
        "Please try again or contact Paul directly."
    )


# -- Email helpers -------------------------------------------------------------

def decode_str(value):
    if value is None: return ""
    decoded, charset = decode_header(value)[0]
    if isinstance(decoded, bytes):
        return decoded.decode(charset or "utf-8", errors="replace")
    return decoded

def get_body(msg):
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            cd = str(part.get("Content-Disposition", ""))
            if ct == "text/plain" and "attachment" not in cd:
                return part.get_payload(decode=True).decode("utf-8", errors="replace")
    else:
        return msg.get_payload(decode=True).decode("utf-8", errors="replace")
    return ""

import re

def format_html_email(text):
    """Convert Amber's markdown-style text into a rich HTML email."""
    lines = text.split("\n")
    html_parts = []
    
    header = """<div style="font-family:-apple-system,BlinkMacSystemFont,'Segoe UI',sans-serif;max-width:700px;margin:0 auto;background:#ffffff;"><div style="background:linear-gradient(135deg,#1a1a2e 0%,#16213e 100%);padding:30px;border-radius:12px 12px 0 0;"><table><tr><td style="width:44px;height:44px;background:#f59e0b;border-radius:50%;text-align:center;font-size:20px;">&#127810;</td><td style="padding-left:12px;"><div style="color:#f59e0b;font-size:18px;font-weight:700;">Agent Amber</div><div style="color:#94a3b8;font-size:12px;">OEP Programme Intelligence</div></td></tr></table></div><div style="padding:24px;background:#f8fafc;">"""
    html_parts.append(header)
    
    in_list = False
    list_type = "ul"
    
    FLAGS = {"JAPAN":"&#127471;&#127477;","SOUTH KOREA":"&#127472;&#127479;","INDIA":"&#127470;&#127475;",
             "BRAZIL":"&#127463;&#127479;","PHILIPPINES":"&#127477;&#127469;","VIETNAM":"&#127483;&#127475;",
             "MEXICO":"&#127474;&#127485;","AUSTRALIA":"&#127462;&#127482;","COLOMBIA":"&#127464;&#127476;"}
    
    def close_list():
        nonlocal in_list
        if in_list:
            html_parts.append(f"</{list_type}>")
            in_list = False
    
    def apply_inline(s):
        s = s.replace("[!]", '<span style="background:#fee2e2;color:#991b1b;padding:2px 8px;border-radius:99px;font-size:11px;font-weight:600;">&#128308; CRITICAL</span>')
        s = s.replace("[~]", '<span style="background:#fef3c7;color:#92400e;padding:2px 8px;border-radius:99px;font-size:11px;font-weight:600;">&#128992; ATTENTION</span>')
        s = s.replace("[-]", '<span style="background:#fef9c3;color:#854d0e;padding:2px 8px;border-radius:99px;font-size:11px;font-weight:600;">&#128993; MONITOR</span>')
        s = s.replace("[ok]", '<span style="background:#dcfce7;color:#166534;padding:2px 8px;border-radius:99px;font-size:11px;font-weight:600;">&#128994; HEALTHY</span>')
        s = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', s)
        s = re.sub(r'\*(.+?)\*', r'<em>\1</em>', s)
        return s
    
    for line in lines:
        line = line.rstrip()
        
        if not line:
            close_list()
            html_parts.append('<div style="height:8px;"></div>')
            continue
        
        if line.startswith("# "):
            close_list()
            html_parts.append(f'<h1 style="font-size:22px;font-weight:700;color:#0f172a;margin:24px 0 12px;padding-bottom:8px;border-bottom:2px solid #e2e8f0;">{apply_inline(line[2:])}</h1>')
        
        elif line.startswith("## "):
            close_list()
            title = line[3:]
            flag = ""
            for country, f in FLAGS.items():
                if country in title.upper():
                    flag = f + " "
                    break
            html_parts.append(f'<div style="background:white;border-radius:8px;padding:16px 20px;margin:12px 0;box-shadow:0 1px 3px rgba(0,0,0,0.08);border-left:4px solid #f59e0b;"><h2 style="font-size:16px;font-weight:700;color:#0f172a;margin:0;">{flag}{apply_inline(title)}</h2></div>')
        
        elif line.startswith("### "):
            close_list()
            html_parts.append(f'<h3 style="font-size:14px;font-weight:600;color:#334155;margin:16px 0 8px;">{apply_inline(line[4:])}</h3>')
        
        elif line.startswith("---"):
            close_list()
            html_parts.append('<hr style="border:none;border-top:1px solid #e2e8f0;margin:16px 0;">')
        
        elif line.startswith("- ") or line.startswith("* "):
            if not in_list:
                html_parts.append('<ul style="margin:8px 0;padding-left:20px;">')
                in_list = True
                list_type = "ul"
            html_parts.append(f'<li style="color:#475569;font-size:14px;margin:4px 0;line-height:1.6;">{apply_inline(line[2:])}</li>')
        
        elif re.match(r'^\d+\.', line):
            if not in_list:
                html_parts.append('<ol style="margin:8px 0;padding-left:20px;">')
                in_list = True
                list_type = "ol"
            item = re.sub(r'^\d+\.\s*', '', line)
            html_parts.append(f'<li style="color:#475569;font-size:14px;margin:4px 0;line-height:1.6;">{apply_inline(item)}</li>')
        
        else:
            close_list()
            html_parts.append(f'<p style="color:#475569;font-size:14px;line-height:1.7;margin:6px 0;">{apply_inline(line)}</p>')
    
    close_list()
    
    footer = """<div style="margin-top:24px;padding:16px;background:#f1f5f9;border-radius:8px;border:1px solid #e2e8f0;text-align:center;"><p style="color:#94a3b8;font-size:12px;margin:0;">Agent Amber &mdash; OEP Programme Intelligence | <a href="mailto:oepagent@oceanenergypathway.org" style="color:#f59e0b;">oepagent@oceanenergypathway.org</a></p></div></div></div>"""
    html_parts.append(footer)
    
    return "\n".join(html_parts)


def send_email(to, subject, body, from_name=None):
    html_body = format_html_email(body)
    data = json.dumps({
        "personalizations": [{"to": [{"email": to}]}],
        "from": {"email": AGENT_EMAIL, "name": from_name or AGENT_NAME},
        "subject": subject,
        "content": [
            {"type": "text/plain", "value": body},
            {"type": "text/html", "value": html_body}
        ]
    }).encode("utf-8")
    req = urllib.request.Request(
        "https://api.sendgrid.com/v3/mail/send",
        data=data,
        headers={"Authorization": f"Bearer {SENDGRID_KEY}", "Content-Type": "application/json"},
        method="POST"
    )
    urllib.request.urlopen(req)
    log.info(f"Email sent to {to} | {subject}")

def is_internal(email_addr):
    return email_addr.lower().endswith(f"@{OEP_DOMAIN}")

# -- Approval flow -------------------------------------------------------------

def send_for_approval(original_from, original_name, original_subject, draft):
    approval_id = f"AMBER-{int(time.time())}"
    pending_approvals[approval_id] = {
        "to": original_from,
        "name": original_name,
        "subject": f"Re: {original_subject}",
        "body": draft,
    }
    body = f"""Hi Paul,

Agent Amber has drafted a reply to an external email. Please APPROVE or REJECT.

------------------------------------
APPROVAL ID: {approval_id}
FROM:        {original_name} <{original_from}>
SUBJECT:     {original_subject}
------------------------------------

DRAFT REPLY:

{draft}

------------------------------------
Reply APPROVE or REJECT.
"""
    send_email(APPROVER_EMAIL, f"[Agent Amber] Approval needed - {approval_id}", body)
    log.info(f"Sent for approval: {approval_id}")

def handle_approval(body, subject):
    match = re.search(r"AMBER-\d+", subject + " " + body)
    if not match:
        return False
    aid = match.group(0)
    if aid not in pending_approvals:
        log.warning(f"Unknown approval ID: {aid}")
        return True
    if "APPROVE" in body.upper():
        item = pending_approvals.pop(aid)
        send_email(item["to"], item["subject"], item["body"])
        send_email(APPROVER_EMAIL, f"[Agent Amber] Sent [sent] - {aid}",
                   f"Reply sent to {item['name']} <{item['to']}>.")
        log.info(f"Approved and sent: {aid}")
    elif "REJECT" in body.upper():
        pending_approvals.pop(aid, None)
        log.info(f"Rejected: {aid}")
    return True

# -- Scheduled tasks -----------------------------------------------------------

def should_run(key, interval_hours):
    last = recall("schedule", key)
    if not last:
        return True
    try:
        last_dt = datetime.fromisoformat(last)
        return (datetime.now(timezone.utc) - last_dt).total_seconds() > interval_hours * 3600
    except:
        return True


def run_strategic_thinking():
    """Amber's daily deep think - synthesises all sources and generates proactive strategic insights for Paul."""
    log.info("Running strategic thinking session...")

    # Max 2 strategic insights per day
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    count_key = f"strategic_count_{today}"
    count = int(recall("strategic", count_key) or "0")
    if count >= 2:
        log.info("Strategic insight limit reached for today (2/2)")
        return

    # Gather context from all sources
    context_parts = []

    # 1. Monday board summary
    try:
        boards_raw = recall("boards", "known_boards") or "{}"
        boards = json.loads(boards_raw)
        if boards:
            context_parts.append("MONDAY.COM BOARDS:")
            for board_id, board_name in list(boards.items())[:8]:
                if is_board_allowed(board_name):
                    context_parts.append(f"  - {board_name} (ID: {board_id})")
    except:
        pass

    # 2. Recent market intel
    try:
        intel = get_recent_intel(days=7)
        if intel:
            context_parts.append("\nRECENT MARKET INTELLIGENCE (last 7 days):")
            for row in intel[:10]:
                country, headline, summary, source, ts = row
                context_parts.append(f"  - {country}: {headline}")
    except:
        pass

    # 3. Recent Slack digest
    try:
        oldest = str((datetime.now(timezone.utc) - timedelta(hours=48)).timestamp())
        auth = slack_get("auth.test", {})
        bot_user_id = auth.get("user_id", "") if auth else ""
        users_result = slack_get("users.list", {"limit": 200})
        user_map = {}
        if users_result and users_result.get("ok"):
            for u in users_result.get("members", []):
                user_map[u["id"]] = u.get("real_name") or u.get("name", u["id"])
        channels_result = slack_get("conversations.list", {
            "types": "public_channel", "exclude_archived": "true", "limit": 100
        })
        if channels_result and channels_result.get("ok"):
            channels = [c for c in channels_result.get("channels", []) if c.get("is_member")]
            slack_msgs = []
            for ch in channels[:8]:
                result = slack_get("conversations.history", {
                    "channel": ch["id"], "oldest": oldest, "limit": 20
                })
                if result and result.get("ok"):
                    for msg in result.get("messages", []):
                        if msg.get("bot_id") or not msg.get("user"):
                            continue
                        uid = msg.get("user", "")
                        if uid == bot_user_id:
                            continue
                        name = user_map.get(uid, uid)
                        text = msg.get("text", "").strip()[:200]
                        if text:
                            slack_msgs.append(f"[#{ch['name']}] {name}: {text}")
            if slack_msgs:
                context_parts.append("\nRECENT SLACK ACTIVITY (last 48h):")
                context_parts.extend(slack_msgs[:40])
    except Exception as e:
        log.error(f"Strategic think - Slack error: {e}")

    # 4. Previously flagged insights (to avoid repeating)
    already_flagged = recall("strategic", "flagged_insights") or ""

    if not context_parts:
        log.info("Not enough context for strategic thinking session")
        return

    context_text = "\n".join(context_parts)

    prompt = f"""You are Agent Amber - a strategic intelligence assistant for Ocean Energy Pathway (OEP), a consultancy accelerating offshore wind development across emerging markets.

You have just completed your daily synthesis of OEP's operational and market intelligence. Your job is to think like a sharp strategy consultant and identify ONE insight that Paul (CEO) genuinely needs to know about.

Think across two lenses simultaneously:
- STRATEGIC OPPORTUNITIES: new markets, untapped partnerships, funding angles, competitive positioning, connections between workstreams nobody has spotted
- OPERATIONAL INTELLIGENCE: patterns in project delays, resource bottlenecks, team dynamics affecting delivery, systemic issues hiding in the data

Here is your full context:
{context_text}

Previously flagged insights (do not repeat): {already_flagged[:800]}

Now think carefully. Ask yourself:
- What pattern is emerging across multiple countries or projects?
- What opportunity exists that OEP hasn't formally pursued?
- What risk is building quietly that Paul might not see yet?
- What connection between two seemingly separate things is worth exploring?
- What would a smart external advisor notice that the team is too close to see?

Only flag something if it is GENUINELY interesting and actionable. Do not fabricate insights - only work from the data above. If nothing stands out, say so.

Respond in JSON only:
{{
  "worth_flagging": true/false,
  "type": "opportunity" or "operational",
  "headline": "One sharp sentence - what you noticed",
  "context": "2-3 sentences explaining the evidence and background",
  "why_it_matters": "1-2 sentences on the strategic or operational significance",
  "options": ["Option 1", "Option 2", "Option 3"],
  "recommendation": "Your single recommended next step - specific and actionable",
  "confidence": "high/medium/low"
}}"""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=800,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = response.content[0].text.strip()
        raw = re.sub(r"^```json|^```|```$", "", raw, flags=re.MULTILINE).strip()
        result = json.loads(raw)

        if not result.get("worth_flagging"):
            log.info("Strategic thinking: nothing worth flagging today")
            return

        headline = result.get("headline", "")
        context_str = result.get("context", "")
        why = result.get("why_it_matters", "")
        options = result.get("options", [])
        recommendation = result.get("recommendation", "")
        insight_type = result.get("type", "opportunity").upper()
        confidence = result.get("confidence", "medium")

        # Avoid repeating similar insights
        if headline[:60] in already_flagged:
            log.info("Strategic thinking: similar insight already flagged, skipping")
            return

        # Format options as numbered list
        options_text = "\n".join(f"{i+1}. {opt}" for i, opt in enumerate(options))

        # Confidence badge
        conf_label = {"high": "High confidence", "medium": "Medium confidence", "low": "Exploratory"}.get(confidence, "")

        body = f"""Hi Paul,

Here's something I've been thinking about - [{insight_type}] {conf_label}

-------------------------------
{headline}
-------------------------------

**What I noticed**
{context_str}

**Why it matters**
{why}

**Your options**
{options_text}

**What I'd recommend**
{recommendation}

-------------------------------
Reply to dig deeper or ask me to explore any of the options above.

- Amber"""

        subject = f"[Amber] {headline}"
        send_email(PAUL_EMAIL, subject, body)
        log.info(f"Strategic insight sent: {headline}")

        # Update counters and memory
        remember("strategic", count_key, str(count + 1))
        new_flagged = (already_flagged + " | " + headline)[-2000:]
        remember("strategic", "flagged_insights", new_flagged)

    except Exception as e:
        log.error(f"Strategic thinking error: {e}")

def run_daily_search():
    log.info("Running daily market intelligence search...")
    markets = [
        "Japan offshore wind energy policy regulation 2026",
        "South Korea offshore wind energy auction Special Act 2026",
        "India offshore wind energy tender Tamil Nadu Gujarat 2026",
        "Brazil offshore wind energy regulation auction COP30 2026",
        "Philippines offshore wind energy DENR auction 2026",
        "Vietnam offshore wind energy MOIT accelerator 2026",
        "Mexico offshore wind energy policy regulation 2026",
        "Australia offshore wind energy auction Victoria 2026",
        "Colombia offshore wind energy auction ANH 2026",
    ]
    for query in markets:
        try:
            execute_tool("web_search", {"query": query, "context": "daily market intelligence"})
            time.sleep(2)
        except Exception as e:
            log.error(f"Search failed for {query}: {e}")
    remember("schedule", "last_daily_search", datetime.now(timezone.utc).isoformat())
    log.info("Daily search complete")

def run_weekly_briefing():
    log.info("Generating weekly Monday briefing for Paul...")
    briefing = run_agentic_loop(
        "Generate the weekly programme intelligence briefing. Cover: "
        "1) Health of all 9 country boards with specific items needing attention (skip completed projects), "
        "2) Market intelligence from this week's web searches, "
        "3) Cross-reference flags between master board and country board activity, "
        "4) Persistent issues that have been flagged 3+ times, "
        "5) Board improvement recommendations, "
        "6) Top 5 priority actions for this week. "
        "Be specific - name items, people, dates.",
        "Agent Amber",
        "amber@internal"
    )
    today = datetime.now().strftime("%d %b %Y")
    send_email(APPROVER_EMAIL, f"[Agent Amber] Weekly Programme Briefing - {today}", briefing)
    remember("schedule", "last_weekly_briefing", datetime.now(timezone.utc).isoformat())
    log.info("Weekly briefing sent")

def run_scheduled_tasks():
    now_utc = datetime.now(timezone.utc)
    # Convert to UK local time (BST = UTC+1 in summer, GMT = UTC+0 in winter)
    # Python's datetime handles DST automatically via astimezone
    try:
        import zoneinfo
        uk_tz = zoneinfo.ZoneInfo("Europe/London")
    except ImportError:
        import pytz
        uk_tz = pytz.timezone("Europe/London")
    now = now_utc.astimezone(uk_tz)
    uk_hour = now.hour

    if should_run("last_daily_search", 24) and uk_hour >= 7:
        run_daily_search()
    # Monday briefing at 8am UK time
    if now.weekday() == 0 and uk_hour >= 8 and uk_hour < 10 and should_run("last_weekly_briefing", 144):
        run_weekly_briefing()
    # Inbox summary at 9am and 3pm UK time
    if uk_hour in (9, 15) and should_run("last_inbox_summary", 5):
        run_inbox_summary()
    # Strategic thinking session daily at 8am UK time
    if uk_hour >= 8 and should_run("last_strategic_think", 20):
        run_strategic_thinking()
        remember("schedule", "last_strategic_think", datetime.now(timezone.utc).isoformat())
    # Proactive intelligence scan - keep running on UTC cycle
    # (already handled in main loop counter)

# -- Main inbox loop -----------------------------------------------------------

def check_inbox():
    try:
        mail = imaplib.IMAP4_SSL(IMAP_SERVER)
        mail.login(AGENT_EMAIL, AGENT_PASSWORD)
        mail.select("inbox")
        _, data = mail.search(None, "UNSEEN")
        ids = data[0].split()
        if not ids:
            log.info("No new emails.")
            mail.logout()
            return
        log.info(f"Found {len(ids)} new email(s).")
        for eid in ids:
            _, msg_data = mail.fetch(eid, "(RFC822)")
            msg = email.message_from_bytes(msg_data[0][1])
            sender_full    = decode_str(msg.get("From", ""))
            sender_subject = decode_str(msg.get("Subject", "(no subject)"))
            body           = get_body(msg)
            m = re.search(r"<(.+?)>", sender_full)
            sender_email = m.group(1) if m else sender_full.strip()
            sender_name  = sender_full.split("<")[0].strip().strip('"') or sender_email
            log.info(f"Email from: {sender_email} | {sender_subject}")
            mail.store(eid, "+FLAGS", "\\Seen")

            # Skip system emails
            skip = ["google.com", "googlemail.com", "mailer-daemon", "sendgrid", "postmaster", "noreply", "no-reply"]
            if any(d in sender_email.lower() for d in skip):
                log.info(f"Skipping system email from {sender_email}")
                continue

            # Handle Paul's approvals
            if sender_email.lower() == APPROVER_EMAIL.lower():
                if handle_approval(body, sender_subject):
                    continue

            # Only respond to authorised senders
            AUTHORISED_SENDERS = [
                "paul@oceanenergypathway.org",
            ]
            if sender_email.lower() not in [s.lower() for s in AUTHORISED_SENDERS]:
                log.info(f"Ignoring unauthorised email from {sender_email}")
                continue

            # Run agentic loop
            try:
                response = run_agentic_loop(body, sender_name, sender_email)
                save_interaction(sender_email, sender_name, body[:1000], response)
                
                # Internal = send directly, external = approval
                # (sender is already verified as internal above, but keeping logic clear)
                if is_internal(sender_email):
                    for attempt in range(3):
                        try:
                            send_email(sender_email, f"Re: {sender_subject}", response)
                            log.info(f"Reply sent directly to {sender_email}")
                            break
                        except Exception as e:
                            log.warning(f"Send attempt {attempt+1} failed: {e}")
                            time.sleep(5)
                else:
                    send_for_approval(sender_email, sender_name, sender_subject, response)
                    
            except Exception as e:
                log.error(f"Agentic loop error for {sender_email}: {e}")
                log.error(traceback.format_exc())
                try:
                    send_email(APPROVER_EMAIL,
                               f"[Agent Amber] Error processing email from {sender_email}",
                               f"Error: {e}\n\nOriginal email:\n{body[:500]}")
                except:
                    pass
        mail.logout()
    except Exception as e:
        log.error(f"Inbox check failed: {e}")

# -- Entry point ---------------------------------------------------------------


# -- Gmail inbox reader (Paul only) --------------------------------------------

def get_gmail_access_token():
    """Get a fresh Gmail access token using the existing Google OAuth refresh token."""
    if not GOOGLE_CLIENT_ID or not GOOGLE_REFRESH_TOKEN:
        return None
    try:
        data = urllib.parse.urlencode({
            "client_id": GOOGLE_CLIENT_ID,
            "client_secret": GOOGLE_CLIENT_SECRET,
            "refresh_token": GOOGLE_REFRESH_TOKEN,
            "grant_type": "refresh_token"
        }).encode()
        req = urllib.request.Request(
            "https://oauth2.googleapis.com/token",
            data=data,
            headers={"Content-Type": "application/x-www-form-urlencoded"},
            method="POST"
        )
        resp = urllib.request.urlopen(req, timeout=10)
        return json.loads(resp.read()).get("access_token")
    except Exception as e:
        log.error(f"Gmail token refresh failed: {e}")
        return None

def gmail_api(path, token, params=None):
    """Call Gmail API."""
    url = f"https://gmail.googleapis.com/gmail/v1/users/me/{path}"
    if params:
        url += "?" + urllib.parse.urlencode(params)
    req = urllib.request.Request(url, headers={"Authorization": f"Bearer {token}"})
    try:
        resp = urllib.request.urlopen(req, timeout=15)
        return json.loads(resp.read())
    except Exception as e:
        log.error(f"Gmail API error ({path}): {e}")
        return None

def read_paul_inbox(hours=24, max_emails=30):
    """Read Paul's inbox and return a structured summary."""
    token = get_gmail_access_token()
    if not token:
        return "Gmail not configured - cannot read inbox."

    # Search for recent emails
    after_ts = int((datetime.now(timezone.utc) - timedelta(hours=hours)).timestamp())
    result = gmail_api("messages", token, {
        "q": f"in:inbox after:{after_ts}",
        "maxResults": max_emails
    })
    if not result:
        return "Could not access Gmail."

    messages = result.get("messages", [])
    if not messages:
        return f"No emails in inbox in the last {hours} hours."

    emails = []
    for msg in messages[:max_emails]:
        detail = gmail_api(f"messages/{msg['id']}", token, {
            "format": "metadata",
            "metadataHeaders": "From,Subject,Date"
        })
        if not detail:
            continue
        headers = {h["name"]: h["value"] for h in detail.get("payload", {}).get("headers", [])}
        snippet = detail.get("snippet", "")[:200]
        label_ids = detail.get("labelIds", [])
        is_unread = "UNREAD" in label_ids
        emails.append({
            "from": headers.get("From", "Unknown"),
            "subject": headers.get("Subject", "(no subject)"),
            "date": headers.get("Date", ""),
            "snippet": snippet,
            "unread": is_unread,
            "id": msg["id"]
        })

    # Sort unread first
    emails.sort(key=lambda x: (not x["unread"], x["date"]))

    lines = []
    unread_count = sum(1 for e in emails if e["unread"])
    lines.append(f"Inbox summary - last {hours}h - {len(emails)} emails ({unread_count} unread):\n")
    for e in emails:
        status = "[!] UNREAD" if e["unread"] else "[sent] read"
        lines.append(f"{status} | {e['date'][:16]} | From: {e['from']} | {e['subject']}")
        if e["snippet"]:
            lines.append(f"   -> {e['snippet']}")
    return "\n".join(lines)

def run_inbox_summary():
    """Generate and send Paul an inbox summary."""
    log.info("Generating inbox summary for Paul...")
    try:
        inbox_data = read_paul_inbox(hours=24, max_emails=40)
        
        prompt = f"""You are Agent Amber. You have just read Paul's inbox at Ocean Energy Pathway.

Here is the raw inbox data:
{inbox_data}

Write Paul a concise, intelligent inbox summary. Structure it as:

1. **Needs your reply** - emails where Paul hasn't responded and a response looks needed
2. **Needs your attention** - important emails, decisions, or time-sensitive items  
3. **FYI / low priority** - updates, newsletters, things he can skim or ignore
4. **Anything you might have missed** - older unread items worth flagging

Be direct and specific. Name the sender and subject. Add a one-line note on why it matters where relevant.
Keep the whole thing under 400 words. Don't pad it out."""

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=800,
            messages=[{"role": "user", "content": prompt}]
        )
        summary = response.content[0].text.strip()
        now_str = datetime.now(timezone.utc).strftime("%A %d %b, %H:%M UTC")
        send_email(PAUL_EMAIL, f"[Amber] Inbox summary - {now_str}", summary)
        remember("schedule", "last_inbox_summary", datetime.now(timezone.utc).isoformat())
        log.info("Inbox summary sent to Paul")
    except Exception as e:
        log.error(f"Inbox summary failed: {e}")

# -- Slack integration --------------------------------------------------------

def slack_api(method, payload):
    """Call a Slack API method."""
    if not SLACK_BOT_TOKEN:
        return None
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(
        f"https://slack.com/api/{method}",
        data=data,
        headers={
            "Authorization": f"Bearer {SLACK_BOT_TOKEN}",
            "Content-Type": "application/json"
        },
        method="POST"
    )
    try:
        resp = urllib.request.urlopen(req, timeout=10)
        return json.loads(resp.read())
    except Exception as e:
        log.error(f"Slack API error ({method}): {e}")
        return None

def slack_post(channel, text):
    """Post a message to a Slack channel."""
    html = format_html_email(text)
    # Slack uses mrkdwn not HTML - convert key formatting
    mrkdwn = text
    mrkdwn = re.sub(r'\*\*(.+?)\*\*', r'*\1*', mrkdwn)
    result = slack_api("chat.postMessage", {
        "channel": channel,
        "text": mrkdwn,
        "mrkdwn": True
    })
    if result and result.get("ok"):
        log.info(f"Slack message posted to {channel}")
    else:
        log.error(f"Slack post failed: {result}")


def slack_get(method, params):
    """Call a Slack API GET method."""
    if not SLACK_BOT_TOKEN:
        return None
    query = urllib.parse.urlencode(params)
    req = urllib.request.Request(
        f"https://slack.com/api/{method}?{query}",
        headers={"Authorization": f"Bearer {SLACK_BOT_TOKEN}"}
    )
    try:
        resp = urllib.request.urlopen(req, timeout=10)
        return json.loads(resp.read())
    except Exception as e:
        log.error(f"Slack GET error ({method}): {e}")
        return None

def is_question(text):
    """Detect if a message is asking a question."""
    text_lower = text.lower().strip()
    # Remove @mentions
    text_clean = re.sub(r'<@[A-Z0-9]+>', '', text_lower).strip()
    if not text_clean:
        return False
    question_words = ['what', 'when', 'where', 'who', 'why', 'how', 'which',
                      'can ', 'could ', 'should ', 'would ', 'is there', 'are there',
                      'do we', 'does ', 'did ', 'has ', 'have we', 'will ']
    if text_clean.endswith('?'):
        return True
    if any(text_clean.startswith(w) for w in question_words):
        return True
    return False

def get_joined_channels():
    """Get list of channels Amber has been invited to."""
    result = slack_get("conversations.list", {
        "types": "public_channel",
        "exclude_archived": "true",
        "limit": 100
    })
    if not result or not result.get("ok"):
        return []
    return [c["id"] for c in result.get("channels", []) if c.get("is_member")]

def scan_slack_channels():
    """Scan channels for unanswered questions and respond."""
    if not SLACK_BOT_TOKEN:
        return
    
    # Track which messages we've already responded to
    responded_key = "slack_responded_ts"
    responded_raw = recall("slack", responded_key) or "{}"
    try:
        responded = json.loads(responded_raw)
    except:
        responded = {}

    # Get bot's own user ID so we don't respond to ourselves
    auth = slack_get("auth.test", {})
    bot_user_id = auth.get("user_id", "") if auth else ""

    channels = get_joined_channels()
    log.info(f"Scanning {len(channels)} Slack channels for questions")

    # Only look at messages from the last 10 minutes
    oldest = str((datetime.now(timezone.utc) - timedelta(minutes=10)).timestamp())

    for channel_id in channels:
        try:
            result = slack_get("conversations.history", {
                "channel": channel_id,
                "oldest": oldest,
                "limit": 20
            })
            if not result or not result.get("ok"):
                continue

            messages = result.get("messages", [])
            for msg in messages:
                ts = msg.get("ts", "")
                text = msg.get("text", "")
                user = msg.get("user", "")

                # Skip our own messages, bots, already responded
                if user == bot_user_id:
                    continue
                if msg.get("bot_id"):
                    continue
                if responded.get(ts):
                    continue
                # Skip if already has a reply from Amber in thread
                if msg.get("reply_count", 0) > 0:
                    responded[ts] = True
                    continue

                if is_question(text):
                    log.info(f"Question detected in channel {channel_id}: {text[:80]}")
                    responded[ts] = True

                    def handle_question(ch=channel_id, t=text, u=user, thread_ts=ts):
                        response = run_agentic_loop(
                            t,
                            f"Slack user {u}",
                            f"slack_{u}@{OEP_DOMAIN}"
                        )
                        slack_mrkdwn = re.sub(r'\*\*(.+?)\*\*', r'*\1*', response)
                        # Reply in thread
                        slack_api("chat.postMessage", {
                            "channel": ch,
                            "text": slack_mrkdwn,
                            "thread_ts": thread_ts,
                            "mrkdwn": True
                        })

                    t = threading.Thread(target=handle_question)
                    t.daemon = True
                    t.start()

        except Exception as e:
            log.error(f"Error scanning channel {channel_id}: {e}")

    # Save responded set (keep last 500 entries to avoid bloat)
    if len(responded) > 500:
        responded = dict(list(responded.items())[-500:])
    remember("slack", responded_key, json.dumps(responded))

def verify_slack_signature(body, timestamp, signature):
    """Verify the request came from Slack."""
    if not SLACK_SIGNING_SECRET:
        return True
    sig_basestring = f"v0:{timestamp}:{body}"
    computed = "v0=" + hmac.new(
        SLACK_SIGNING_SECRET.encode(),
        sig_basestring.encode(),
        hashlib.sha256
    ).hexdigest()  # hmac.new is correct Python stdlib
    return hmac.compare_digest(computed, signature)

class SlackHandler(BaseHTTPRequestHandler):
    """HTTP handler for Slack events."""
    
    def do_POST(self):
        content_length = int(self.headers.get("Content-Length", 0))
        body_bytes = self.rfile.read(content_length)
        body = body_bytes.decode("utf-8")
        
        try:
            payload = json.loads(body)
            
            # URL verification - skip signature check, respond immediately
            if payload.get("type") == "url_verification":
                challenge = payload["challenge"]
                response_body = json.dumps({"challenge": challenge}).encode("utf-8")
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.send_header("Content-Length", str(len(response_body)))
                self.end_headers()
                self.wfile.write(response_body)
                log.info("Slack URL verification challenge answered")
                return
            
            # Verify signature for all real events
            timestamp = self.headers.get("X-Slack-Request-Timestamp", "")
            signature = self.headers.get("X-Slack-Signature", "")
            if not verify_slack_signature(body, timestamp, signature):
                self.send_response(403)
                self.end_headers()
                return
            
            self.send_response(200)
            self.end_headers()
            
            # Handle events
            event = payload.get("event", {})
            event_type = event.get("type")
            
            # Respond to @mentions
            if event_type == "app_mention":
                channel = event.get("channel")
                user = event.get("user")
                text = event.get("text", "")
                
                # Remove the @Amber mention from the text
                text = re.sub(r'<@[A-Z0-9]+>', '', text).strip()
                
                log.info(f"Slack mention in {channel}: {text[:100]}")
                
                # Run in background thread so we return 200 quickly
                def handle_mention():
                    response = run_agentic_loop(
                        text,
                        f"Slack user {user}",
                        f"slack_{user}@{OEP_DOMAIN}"
                    )
                    # Convert to Slack markdown
                    slack_response = re.sub(r'\*\*(.+?)\*\*', r'*\1*', response)
                    slack_post(channel, slack_response)
                
                thread = threading.Thread(target=handle_mention)
                thread.daemon = True
                thread.start()
                
        except Exception as e:
            log.error(f"Slack event handling error: {e}")
    
    def log_message(self, *args):
        pass  # Suppress default HTTP logging

def start_slack_server():
    """Start the HTTP server for Slack events on port 8080."""
    if not SLACK_BOT_TOKEN:
        log.info("Slack not configured - skipping")
        return
    try:
        server = HTTPServer(("0.0.0.0", 8080), SlackHandler)
        thread = threading.Thread(target=server.serve_forever)
        thread.daemon = True
        thread.start()
        log.info("Slack event server started on port 8080")
    except Exception as e:
        log.error(f"Slack server failed to start: {e}")



def run_proactive_intelligence():
    """Scan all sources and decide if anything is worth flagging to Paul proactively."""
    if not SLACK_BOT_TOKEN:
        return

    # Max 3 proactive emails per day
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    count_key = f"proactive_count_{today}"
    count = int(recall("proactive", count_key) or "0")
    if count >= 3:
        log.info("Proactive email limit reached for today (3/3)")
        return

    log.info("Running proactive intelligence scan...")

    # Gather Slack context from last 6 hours across all channels
    hours = 6
    oldest = str((datetime.now(timezone.utc) - timedelta(hours=hours)).timestamp())

    auth = slack_get("auth.test", {})
    bot_user_id = auth.get("user_id", "") if auth else ""

    users_result = slack_get("users.list", {"limit": 200})
    user_map = {}
    if users_result and users_result.get("ok"):
        for u in users_result.get("members", []):
            user_map[u["id"]] = u.get("real_name") or u.get("name", u["id"])

    channels_result = slack_get("conversations.list", {
        "types": "public_channel", "exclude_archived": "true", "limit": 100
    })
    if not channels_result or not channels_result.get("ok"):
        return

    channels = [c for c in channels_result.get("channels", []) if c.get("is_member")]
    slack_digest = []
    for ch in channels:
        result = slack_get("conversations.history", {
            "channel": ch["id"], "oldest": oldest, "limit": 30
        })
        if not result or not result.get("ok"):
            continue
        for msg in result.get("messages", []):
            if msg.get("bot_id") or msg.get("subtype"):
                continue
            uid = msg.get("user", "")
            if uid == bot_user_id:
                continue
            name = user_map.get(uid, uid)
            text = msg.get("text", "").strip()
            if text:
                slack_digest.append(f"[#{ch['name']}] {name}: {text[:200]}")

    if not slack_digest:
        log.info("No Slack activity in last 6 hours - skipping proactive scan")
        return

    # Get already-flagged insights to avoid repeating
    already_flagged = recall("proactive", "flagged_summaries") or ""

    slack_text = "\n".join(slack_digest[:60])  # cap to avoid token overload

    prompt = f"""You are Agent Amber, an intelligent assistant for Ocean Energy Pathway (OEP).

Below is a digest of Slack conversations from the last 6 hours across OEP channels.

Your job: read this carefully and decide if there is anything genuinely worth flagging proactively to Paul (the CEO). 

Flag things like:
- A team member seems blocked, frustrated, or struggling to get traction
- A relationship or collaboration seems strained
- A project or country workstream seems at risk
- An opportunity is being discussed that Paul should be aware of
- A decision is being delayed that needs his input
- Any pattern suggesting something needs leadership attention

Do NOT flag:
- Normal day-to-day work updates
- Things already resolved in the conversation
- Minor admin matters

Recent Slack activity:
{slack_text}

Previously flagged (don't repeat these): {already_flagged[:500]}

Respond in JSON only:
{{
  "worth_flagging": true/false,
  "insight": "One clear sentence describing what you spotted",
  "detail": "2-3 sentences of context and why Paul should care",
  "suggested_action": "One concrete suggestion for what Paul might do"
}}"""

    try:
        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=500,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = response.content[0].text.strip()
        # Strip markdown fences if present
        raw = re.sub(r"^```json|^```|```$", "", raw, flags=re.MULTILINE).strip()
        result = json.loads(raw)

        if not result.get("worth_flagging"):
            log.info("Proactive scan: nothing worth flagging")
            return

        insight = result.get("insight", "")
        detail = result.get("detail", "")
        action = result.get("suggested_action", "")

        # Check we haven't flagged something very similar recently
        if insight[:50] in already_flagged:
            log.info(f"Proactive scan: already flagged similar insight, skipping")
            return

        # Send proactive email to Paul
        subject = f"[Amber] {insight}"
        body = f"""Hi Paul,

I wanted to flag something I noticed from recent Slack activity:

**{insight}**

{detail}

**Suggested action:** {action}

This is an automated insight from my monitoring of OEP channels. Reply if you'd like me to dig deeper.

- Amber"""

        send_email(APPROVER_EMAIL, subject, body)
        log.info(f"Proactive email sent: {insight}")

        # Update count and memory
        remember("proactive", count_key, str(count + 1))
        new_flagged = (already_flagged + " | " + insight)[-1000:]
        remember("proactive", "flagged_summaries", new_flagged)

    except Exception as e:
        log.error(f"Proactive intelligence error: {e}")

if __name__ == "__main__":
    log.info(f"Agent Amber v36 starting - polling every {POLL_INTERVAL}s")
    log.info(f"Watching:  {AGENT_EMAIL}")
    log.info(f"Approver:  {APPROVER_EMAIL}")
    log.info(f"Domain:    @{OEP_DOMAIN}")
    log.info(f"Drive:     {'connected' if GOOGLE_CLIENT_ID else 'not configured'}")

    init_db()
    start_slack_server()

    slack_scan_counter = 0
    proactive_counter = 0
    while True:
        try:
            run_scheduled_tasks()
            check_inbox()
            # Scan Slack for questions every ~5 minutes
            slack_scan_counter += 1
            if slack_scan_counter >= 3:
                scan_slack_channels()
                slack_scan_counter = 0
            # Proactive intelligence scan every ~4 hours
            proactive_counter += 1
            if proactive_counter >= 120:
                run_proactive_intelligence()
                proactive_counter = 0
        except Exception as e:
            log.error(f"Main loop error: {e}")
        time.sleep(POLL_INTERVAL)
